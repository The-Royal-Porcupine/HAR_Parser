import os
import pandas as pd
from pptx import Presentation
from pd2ppt import df_to_table
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from datetime import datetime

import HAR_Analysis as har
import utils
import objects

# Some Technical Functions for PPTX


def format_slide_table(slide, new_table, title_font, row_font, col_idx=None, col_wid=None):
    # Setting font size for header
    for cell in new_table.rows[0].cells:
        for par in cell.text_frame.paragraphs:
            par.font.size = Pt(title_font)
    # Setting font size for lines
    for r in range(1, len(new_table.rows)):
        for cell in new_table.rows[r].cells:
            for par in cell.text_frame.paragraphs:
                par.font.size = Pt(row_font)
    # Setting vertical alignment for all
    for cell in new_table.iter_cells():
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    # Setting column width
    total_table_width = 0.0
    for cc in new_table.columns:
        width_cm = round(float(cc.width / 360000), 3)
        total_table_width += width_cm
    if col_idx and col_wid:
        subset_width = 0.0
        for ci in col_idx:
            new_table.columns[ci].width = Cm(col_wid[col_idx.index(ci)])
            subset_width += Cm(col_wid[col_idx.index(ci)])
        rest_width = Cm(total_table_width) - subset_width
        col_num = len(new_table.columns)
        for i in range(0, col_num):
            if i in col_idx:
                continue
            else:
                new_table.columns[i].width = round(rest_width / (col_num - len(col_idx)))


def create_slide(layout=3, title_text=None, subtitle_text=None, presentation=None):
    presa = presentation
    slide = presa.slides.add_slide(presa.slide_layouts[layout])
    if slide.shapes.title:
        title_placeholder = slide.shapes.title
        title_placeholder.text = title_text
    for shape in slide.placeholders:
        if 'Subtitle' in shape.name:
            shape.text = subtitle_text
    #print('\nNew slide with following placeholders created:')
    #for shape in slide.placeholders:
        #print(f'{shape.placeholder_format.idx}, {shape.name}')
    return slide


def create_presentation(path_to_har_file, path_to_pptx):

    ########## 1.1 Get Objects

    analysis = har.StoryAnalyzer(path_to_har_file)
    analysis.read_har_file()

    # 1.1.1 Product and Story

    products = analysis.get_product_info()
    story = analysis.get_story_info_2()

    # 1.1.2. All widgets existing in Story (from ContentLib) + Calculation Entities

    widgets_content_lib = analysis.get_widgets_contentlib_info()[0]
    calc_entities = analysis.get_widgets_contentlib_info()[1]

    # 1.1.3 Get Calls Info

    calls = analysis.get_calls_info(calc_entities)

    # 1.1.4 Get all widgets that have been loaded (Events)

    events_widgets = analysis.get_widgets_info()

    # 1.1.5 Get all widgets with MDS Requests

    # test_calls + test_widgets

    # 1.1.6 Action Summary

    # events_widgets

    # 1.1.7 Widgets + MDS Execution Timeline

    # calls + events_widgets

    # 1.1.8 All Widgets Timeline

    # events_widgets

    # 1.1.9 MDS Request for Widget Info

    presentation = Presentation(os.getcwd() + '/SAP_Template_Presentation.pptx')

    # 1 Slide - Title

    if story:

        story_frame = utils.get_story_frame(story)
        title_text = story_frame.T['Story Name'][0]
        story_title = f'Story Analysis ({title_text})'
        slide1 = create_slide(layout=0, title_text=story_title, subtitle_text=f'by {os.getlogin()}', presentation=presentation)
        image = slide1.placeholders[12]
        image.insert_picture(os.getcwd() + '/banner.png')

    # 2 Slide - Total Story Info And Widgets Total Count - Two Tables

    if products:

        product_frame = utils.get_product_frame(products)
        story_product_frame = pd.concat((product_frame, story_frame))
        title = 'General Story Info And Widgets Total Count'
        slide2 = create_slide(layout=1, title_text=title, presentation=presentation)
        table1 = slide2.placeholders[10]
        df = story_product_frame.reset_index()
        df = df.rename(columns=df.iloc[0])
        df_to_table(slide2, df, left=table1.left, top=table1.top, width=table1.width, height=table1.height)

    if widgets_content_lib:

        widgets_content_lib_frame = utils.get_widgets_from_contentlib_frame(widgets_content_lib)
        table2 = slide2.placeholders[11]
        widgets_frame = widgets_content_lib_frame[['Story Name', 'Page Title', 'Widget ID', 'Widget Class']]
        widgets_frame_columns = ['Story Name', 'Page Title', 'Number of Widgets', 'Widget Class']
        widgets_frame.columns = widgets_frame_columns
        widgets_frame = widgets_frame.groupby(['Story Name', 'Page Title', 'Widget Class']).count()
        df_to_table(slide2, widgets_frame.reset_index(), left=table2.left, top=table2.top,
                    width=table2.width, height=table2.height)

    # 3 Slide - Time Graphs - Pie Charts

    if calls:

        story_summary = utils.get_story_summary(calls)
        slide3 = create_slide(layout=4, title_text='Time Graphs in Seconds', presentation=presentation)
        slide3.placeholders[17].text = 'Total and average time graphs are made with wait time excluded'
        #slide3.placeholders[17].text.font.size = Pt(18)
        total_time_chart_shape = slide3.placeholders[22]
        average_time_chart_shape = slide3.placeholders[23]
        biggest_time_chart_shape = slide3.placeholders[24]

        pie_charts_df = story_summary
        wait_index = pie_charts_df.index.isin(['wait'])
        # Prepare Data
        total_time_chart_data = CategoryChartData()
        # Labels
        total_time_chart_data.categories = pie_charts_df[~wait_index].index
        # Data
        total_time_chart_data.add_series('Series 1', pie_charts_df['total_runtime'][~wait_index].map('{:,.2f}'.format))
        total_time_chart_shape = total_time_chart_shape.insert_chart(XL_CHART_TYPE.PIE, total_time_chart_data)
        total_time_chart = total_time_chart_shape.chart
        # Format Chart
        total_time_chart.has_legend = True
        total_time_chart.legend.position = XL_LEGEND_POSITION.CORNER
        total_time_chart.legend.font.size = Pt(9)
        total_time_chart.legend.include_in_layout = False
        total_time_chart.plots[0].has_data_labels = True
        total_time_chart_data_labels = total_time_chart.plots[0].data_labels
        total_time_chart_data_labels.font.size = Pt(9)
        total_time_chart_data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        slide3.placeholders[15].text = 'Total Trace Time with Wait Time Excluded in Seconds'
        #slide3.placeholders[15].text.font.size = Pt(16)

        average_time_chart_data = CategoryChartData()
        average_time_chart_data.categories = pie_charts_df[~wait_index].index
        average_time_chart_data.add_series('Series 1', pie_charts_df['average_runtime'][~wait_index].map('{:,.2f}'.format))
        average_time_chart_shape = average_time_chart_shape.insert_chart(XL_CHART_TYPE.PIE, average_time_chart_data)
        average_time_chart = average_time_chart_shape.chart
        average_time_chart.has_legend = True
        average_time_chart.legend.position = XL_LEGEND_POSITION.CORNER
        average_time_chart.legend.font.size = Pt(9)
        average_time_chart.legend.include_in_layout = False
        average_time_chart.plots[0].has_data_labels = True
        average_time_chart_data_labels = average_time_chart.plots[0].data_labels
        average_time_chart_data_labels.font.size = Pt(9)
        average_time_chart_data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        slide3.placeholders[20].text = 'Average Times with Wait Time Excluded'
        #slide3.placeholders[20].text.font.size = Pt(16)

        biggest_time_chart_data = CategoryChartData()
        biggest_time_chart_data.categories = pie_charts_df.index
        biggest_time_chart_data.add_series('Series 1', pie_charts_df['maximum_runtime'].map('{:,.2f}'.format))
        biggest_time_chart_shape = biggest_time_chart_shape.insert_chart(XL_CHART_TYPE.PIE, biggest_time_chart_data)
        biggest_time_chart = biggest_time_chart_shape.chart
        biggest_time_chart.has_legend = True
        biggest_time_chart.legend.position = XL_LEGEND_POSITION.CORNER
        biggest_time_chart.legend.font.size = Pt(9)
        biggest_time_chart.legend.include_in_layout = False
        biggest_time_chart.plots[0].has_data_labels = True
        biggest_time_chart_data_labels = biggest_time_chart.plots[0].data_labels
        biggest_time_chart_data_labels.font.size = Pt(9)
        biggest_time_chart_data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        slide3.placeholders[21].text = 'The biggest times'
        #slide3.placeholders[21].text.font.size = Pt(16)

    # 4 Slide - Widgets executed when opening page(s)

    if events_widgets:

        df = utils.get_widgets_from_events_frame(events_widgets)[0]
        actions = df['User Action'].unique().tolist()
        for action in actions:
            #events_widgets_df = df.query("'User Action' == @action")
            events_widgets_df = df[df['User Action'] == action].drop(columns=['User Action Start Time'])
            max_lines_on_slide = 15
            slide_number = 0

            while len(events_widgets_df.index) > 0:
                if len(events_widgets_df.index) > max_lines_on_slide:
                    slide_number += 1
                    slide4 = create_slide(layout=2, title_text='Widgets executed during user action: "' + action + '" (' + str(slide_number) + ') ',
                                  presentation=presentation)
                    table = slide4.placeholders[10]
                    slide4_df = df_to_table(slide4, events_widgets_df[0:(max_lines_on_slide - 1)], left=table.left, top=table.top,
                                width=table.width, height=table.height)
                    events_widgets_df = events_widgets_df[max_lines_on_slide:len(events_widgets_df.index)]
                    format_slide_table(slide4, slide4_df.table, title_font=12, row_font=10, col_idx=[0,1,2,3,4,5],
                                       col_wid=[6.8, 3.4, 6.2, 2.3, 4.6, 4.6, 3.2])
                    row_height = slide4_df.table.rows[0].height
                else:
                    slide_number += 1
                    slide4 = create_slide(layout=2, title_text='Widgets executed during user action: "' + action + '" (' + str(slide_number) + ') ',
                                          presentation=presentation)
                    table = slide4.placeholders[10]
                    slide4_df = df_to_table(slide4, events_widgets_df, left=table.left, top=table.top,
                                width=table.width, height=table.height)
                    format_slide_table(slide4, slide4_df.table, title_font=12, row_font=10, col_idx=[0,1,2,3,4,5],
                                       col_wid=[6.8, 3.4, 6.2, 2.3, 4.6, 4.6, 3.2])
                    events_widgets_df = events_widgets_df.iloc[0:0]
                    row_height = slide4_df.table.rows[0].height
                    for row in slide4_df.table.rows:
                        row.height = Cm(0.8)

    # 5 Slide - Widgets with MDS Requests

    if calls and events_widgets:

        df = utils.get_widgets_from_mds_frame(calls, events_widgets)[0]
        actions = df['User Action'].unique().tolist()
        for action in actions:
            mds_widgets_df = df[df['User Action'] == action].drop(columns=['User Action'])
            max_lines_on_slide = 15
            slide_number = 0

            while len(mds_widgets_df.index) > 0:
                if len(mds_widgets_df.index) > max_lines_on_slide:
                    slide_number += 1
                    slide5 = create_slide(layout=2, title_text='Widgets with MDS requests for action "' + action + '" (' + str(slide_number) + ') ',
                                          presentation=presentation)
                    table = slide5.placeholders[10]
                    slide5_df = df_to_table(slide5, mds_widgets_df[0:(max_lines_on_slide - 1)], left=table.left,
                                            top=table.top,
                                            width=table.width, height=table.height)
                    mds_widgets_df = mds_widgets_df[max_lines_on_slide:len(mds_widgets_df.index)]
                    format_slide_table(slide5, slide5_df.table, title_font=12, row_font=10, col_idx=[0, 1, 3, 4],
                                       col_wid=[2.8, 6.2, 4.6, 4.6])
                else:
                    slide_number += 1
                    slide5 = create_slide(layout=2, title_text='Widgets with MDS requests for action "' + action + '" (' + str(slide_number) + ') ',
                                          presentation=presentation)
                    table = slide5.placeholders[10]
                    slide5_df = df_to_table(slide5, mds_widgets_df, left=table.left, top=table.top,
                                            width=table.width, height=table.height)
                    format_slide_table(slide5, slide5_df.table, title_font=12, row_font=10, col_idx=[0, 1, 3, 4],
                                       col_wid=[2.8, 6.2, 4.6, 4.6])
                    mds_widgets_df = mds_widgets_df.iloc[0:0]

    # 6 Slide - Action Summary

    if events_widgets:

        df = utils.get_widgets_from_events_frame(events_widgets)[1]
        table_summary = df.groupby('User Action', as_index=False) \
                          .agg({'User Action Start': lambda x: x.iloc[0], 'Widget Start':'min', 'Widget Finish':'max'}) \
                          .sort_values('User Action Start') \
                          .rename(columns={'Widget Start':'First Widget Start', 'Widget Finish':'Last Widget End'})
        table_summary['Widget Start to End (sec)'] = table_summary['Last Widget End'] - table_summary['First Widget Start']
        table_summary['Widget Start to End (sec)'] = round(table_summary['Widget Start to End (sec)'].dt.total_seconds(), 3)
        table_summary['Action Start to Widget End (sec)'] = table_summary['Last Widget End'] - table_summary['User Action Start']
        table_summary['Action Start to Widget End (sec)'] = round(table_summary['Action Start to Widget End (sec)'].dt.total_seconds(), 3)
        table_summary['User Action Start'] = table_summary['User Action Start'].dt.strftime("%H:%M:%S.%f")
        table_summary['First Widget Start'] = table_summary['First Widget Start'].dt.strftime("%H:%M:%S.%f")
        table_summary['Last Widget End'] = table_summary['Last Widget End'].dt.strftime("%H:%M:%S.%f")

        slide6 = create_slide(layout=2, title_text='Summary of User Actions', presentation=presentation)
        table = slide6.placeholders[10]
        slide6_df = df_to_table(slide6, table_summary, left=table.left,
                                top=table.top,
                                width=table.width)
        format_slide_table(slide6, slide6_df.table, title_font=12, row_font=12, col_idx=[1, 2, 3, 4, 5],
                           col_wid=[5.8, 5.8, 5.8, 3.5, 3.5])
        for row in slide6_df.table.rows:
            row.height = Cm(0.8)

    # 7 Slide - Widget & MDS Execution timeline

    if calls and events_widgets:

        df = utils.get_widgets_from_mds_frame(calls, events_widgets)[1]
        actions = df['User Action'].unique().tolist()
        for action in actions:
            gantt_df = df[df['User Action'] == action]
            slide7 = create_slide(layout=5, title_text='Widget & MDS Execution timeline for action: "' + action + '" ', presentation=presentation)
            gantt_chart_shape = slide7.placeholders[22]
            gantt_chart_data = CategoryChartData()
            gantt_chart_data.categories = gantt_df['MDS Number'].values
            gantt_chart_data.add_series('', gantt_df['Start Mark'], number_format="hh:mm:ss.000")
            gantt_chart_data.add_series('Widget ->', gantt_df['Before MDS Duration'], number_format="hh:mm:ss.000")
            gantt_chart_data.add_series('MDS Request', gantt_df['MDS Duration'], number_format="hh:mm:ss.000")
            gantt_chart_data.add_series('<- Widget', gantt_df['After MDS Duration'], number_format="hh:mm:ss.000")
            gantt_chart_shape = gantt_chart_shape.insert_chart(XL_CHART_TYPE.BAR_STACKED, gantt_chart_data)
            gantt_chart = gantt_chart_shape.chart
            gantt_chart.has_legend = True
            gantt_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            gantt_chart.legend.include_in_layout = False

            category_axis = gantt_chart.category_axis
            category_axis.tick_labels.font.size = Pt(12)

            value_axis = gantt_chart.value_axis
            value_axis.tick_labels.font.size = Pt(12)

            plot = gantt_chart.plots[0]
            plot.gap_width = 50
            plot.series[0].format.fill.background()
            plot.series[1].format.fill.solid()
            plot.series[2].format.fill.solid()
            plot.series[3].format.fill.solid()
            plot.series[0].format.line.fill.background()
            plot.series[1].format.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_4
            plot.series[3].format.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_4
            plot.series[2].format.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_5
            plot.has_data_labels = True
            plot.data_labels.font.size = Pt(10)
            plot.data_labels.number_format = "mm:ss.000"
            plot.series[0].data_labels.font.fill.background()
            plot.series[1].data_labels.font.color.rgb = RGBColor(0x01, 0x01, 0x01)
            plot.series[3].data_labels.font.color.rgb = RGBColor(0x01, 0x01, 0x01)

    # 8 Slide - Gantt Common Widgets timelines (by Action)

    if events_widgets:

        df = utils.get_widgets_from_events_frame(events_widgets)[1]
        actions = df['User Action'].unique().tolist()
        for action in actions:
            gantt_df = df[df['User Action'] == action]
            slide8 = create_slide(layout=5, title_text='All Widgets Timeline for User Action: "' + action + '" ', presentation=presentation)
            gantt_chart_shape = slide8.placeholders[22]
            gantt_chart_data = CategoryChartData()

            gantt_chart_data.categories = gantt_df['Label'].values
            gantt_chart_data.add_series('', gantt_df['Widget Start'], number_format="hh:mm:ss.000")
            gantt_chart_data.add_series('Widget runtime', gantt_df['Widget Duration'], number_format="hh:mm:ss.000")

            gantt_chart_shape = gantt_chart_shape.insert_chart(XL_CHART_TYPE.BAR_STACKED, gantt_chart_data)
            gantt_chart = gantt_chart_shape.chart
            gantt_chart.has_legend = True
            gantt_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            gantt_chart.legend.include_in_layout = False

            category_axis = gantt_chart.category_axis
            category_axis.tick_labels.font.size = Pt(10)
            value_axis = gantt_chart.value_axis
            value_axis.tick_labels.font.size = Pt(9)
            plot = gantt_chart.plots[0]
            plot.gap_width = 50
            plot.series[0].format.fill.background()
            plot.series[0].format.line.fill.background()
            plot.series[0].data_labels.font.fill.background()
            plot.series[1].format.fill.solid()
            #plot.series[1].format.fill.fore_color.rgb = RGBColor(0x9A, 0xC8, 0xFC)
            plot.series[1].format.fill.fore_color.rgb = RGBColor(0, 188, 242)
            plot.has_data_labels = True
            plot.data_labels.font.size = Pt(9)
            plot.data_labels.number_format = "mm:ss.000"

    # 9 Slides with MDS Info and MDS Definition

    if calls:

        index = 1
        for call in calls:
            if isinstance(call, objects.Call_2) and call.get_response_type == 'Batch':
                # First Slide: tables with MDS Details and Views Info + graph
                mds_timings_frame = utils.get_mds_timings_frame(call, index)
                mds_timings_frame = mds_timings_frame.reset_index()
                mds_timings_frame = mds_timings_frame.rename(columns=mds_timings_frame.iloc[0]).drop(index=0)
                mds_views_frame = utils.get_view_info_table(call)
                slide9 = create_slide(layout=6, title_text='MDS Request for widget', presentation=presentation)
                if call.widgets:
                    slide9.placeholders[25].text = f'Widget ID: {call.widgets[0]}'
                timings_table = slide9.placeholders[23]
                views_table = slide9.placeholders[24]
                slide9_timings_df = df_to_table(slide9, mds_timings_frame, left=timings_table.left, top=timings_table.top,
                                                width=timings_table.width, height=timings_table.height)
                slide9_views_df = df_to_table(slide9, mds_views_frame, left=views_table.left, top=views_table.top,
                                              width=views_table.width, height=views_table.height)
                format_slide_table(slide9, slide9_timings_df.table, title_font=14, row_font=12, col_idx=[0], col_wid=[4.5])
                format_slide_table(slide9, slide9_views_df.table, title_font=14, row_font=10, col_idx=[0,2,3,4],
                                   col_wid=[8.6, 4.5, 9, 2.5])
                for row in slide9_views_df.table.rows:
                    row.height = Cm(0.8)
                index += 1

                chart_shape = slide9.placeholders[22]
                chart_data = CategoryChartData()

                chart_data.categories = [ds.ds_name for ds in call.datasources]
                chart_data.add_series('Runtime (s)', [ds.ds_runtime for ds in call.datasources])

                chart_shape = chart_shape.insert_chart(XL_CHART_TYPE.BAR_STACKED, chart_data)
                chart = chart_shape.chart

                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(14)

                value_axis = chart.value_axis
                value_axis.tick_labels.font.size = Pt(10)
                value_axis.axis_title.text_frame.paragraphs[0].text = 'Runtime (s)'
                value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(12)

                plot = chart.plots[0]
                plot.gap_width = 30
                plot.vary_by_categories = False
                plot.series[0].format.fill.solid()
                plot.series[0].format.fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_4

                plot.has_data_labels = True
                plot.data_labels.show_category_name = False
                plot.data_labels.show_value = False
                plot.data_labels.font.size = Pt(12)

                dpl_ind = 0
                for p in plot.series[0].points:
                    p.data_label.text_frame.auto_size = True
                    p.data_label.text_frame.word_wrap = False
                    p.data_label.text_frame.paragraphs[0].add_run()
                    p.data_label.text_frame.paragraphs[0].runs[0].text = [ds.ds_name for ds in call.datasources][dpl_ind]
                    p.data_label.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
                    p.data_label.position = XL_DATA_LABEL_POSITION.INSIDE_BASE
                    p.data_label.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    dpl_ind += 1

                # Second Slide(s): MDS Details: Dimensions, Calc Entities
                for ds in call.datasources:
                    slide10 = create_slide(layout=7, title_text=f'Definition for MDS "{ds.ds_name}"', presentation=presentation)
                    slide10.placeholders[12].text = f'View: {ds.view}'

                    text_shape = slide10.placeholders[11]
                    text_frame = text_shape.text_frame
                    text_frame.fit_text(max_size=18)

                    p1 = text_frame.paragraphs[0]
                    p1.text = f'MDS Runtime: {ds.ds_runtime} (s)'
                    font = p1.runs[0].font
                    font.name = 'Calibri'

                    p2 = text_frame.add_paragraph()
                    p2.text = f'MDS Dimensions:'
                    font = p2.runs[0].font
                    font.name = 'Calibri'
                    font.underline = True

                    for k in ds.dimensions.keys():
                        p = text_frame.add_paragraph()
                        p.level = 1
                        p.space_after = Pt(1)
                        p.space_before = Pt(1)
                        p.text = f'{k}'
                        p.runs[0].font.size = Pt(12)
                        if ds.dimensions[k]:
                            pt = text_frame.add_paragraph()
                            pt.text = f'Formula Members in {k}:'
                            font = pt.runs[0].font
                            font.name = 'Calibri'
                            font.underline = True
                            for item in ds.dimensions[k]:
                                pi = text_frame.add_paragraph()
                                pi.level = 1
                                pi.space_after = Pt(1)
                                pi.space_before = Pt(1)
                                pi.text = f'{item}'
                                pi.runs[0].font.size = Pt(12)

    # 10 Appendix 1 - All Widgets Existing in Story

    if widgets_content_lib:

        df = widgets_content_lib_frame
        pages = df['Page Title'].unique().tolist()
        for page in pages:
            contlib_widgets_df = df[df['Page Title'] == page].drop(columns=['Story Name'])
            max_lines_on_slide = 15
            slide_number = 0
            while len(contlib_widgets_df.index) > 0:
                if len(contlib_widgets_df.index) > max_lines_on_slide:
                    slide_number += 1
                    slide11 = create_slide(layout=2,
                                          title_text='Appendix 1: All widgets used on page: "' + page + '" (' + str(
                                              slide_number) + ') ',
                                          presentation=presentation)
                    table = slide11.placeholders[10]
                    slide11_df = df_to_table(slide11, contlib_widgets_df[0:(max_lines_on_slide - 1)], left=table.left,
                                            top=table.top,
                                            width=table.width, height=table.height)
                    contlib_widgets_df = contlib_widgets_df[max_lines_on_slide:len(contlib_widgets_df.index)]
                    format_slide_table(slide11, slide11_df.table, title_font=12, row_font=10, col_idx=[2, 3],
                                       col_wid=[9.0, 10.0])
                    row_height = slide11_df.table.rows[0].height
                else:
                    slide_number += 1
                    slide11 = create_slide(layout=2,
                                          title_text='Appendix 1: All widgets used on page: "' + page + '" (' + str(
                                              slide_number) + ') ',
                                          presentation=presentation)
                    table = slide11.placeholders[10]
                    slide11_df = df_to_table(slide11, contlib_widgets_df, left=table.left, top=table.top,
                                            width=table.width, height=table.height)
                    format_slide_table(slide11, slide11_df.table, title_font=12, row_font=10, col_idx=[2, 3],
                                       col_wid=[9.0, 10.0])
                    contlib_widgets_df = contlib_widgets_df.iloc[0:0]
                    row_height = slide11_df.table.rows[0].height
                    for row in slide11_df.table.rows:
                        row.height = Cm(0.8)

    # 11 Appendix 2 - All Calculation Entities in Story

    if calc_entities:

        df = utils.get_calc_entities_frame(calc_entities)
        max_lines_on_slide = 15
        slide_number = 0
        while len(df.index) > 0:
            if len(df.index) > max_lines_on_slide:
                slide_number += 1
                slide12 = create_slide(layout=2,
                                      title_text='Appendix 2: All calculation entities used in story (' + str(
                                          slide_number) + ') ',
                                      presentation=presentation)
                table = slide12.placeholders[10]
                slide12_df = df_to_table(slide12, df[0:(max_lines_on_slide - 1)], left=table.left,
                                        top=table.top,
                                        width=table.width, height=table.height)
                df = df[max_lines_on_slide:len(df.index)]
                format_slide_table(slide12, slide12_df.table, title_font=12, row_font=10, col_idx=[0],
                                   col_wid=[10.0])
                row_height = slide12_df.table.rows[0].height
            else:
                slide_number += 1
                slide12 = create_slide(layout=2,
                                      title_text='Appendix 2: All calculation entities used in story (' + str(
                                          slide_number) + ') ',
                                      presentation=presentation)
                table = slide12.placeholders[10]
                slide12_df = df_to_table(slide12, df, left=table.left, top=table.top,
                                        width=table.width, height=table.height)
                format_slide_table(slide12, slide12_df.table, title_font=12, row_font=10, col_idx=[0],
                                   col_wid=[10.0])
                df = df.iloc[0:0]
                row_height = slide12_df.table.rows[0].height
                for row in slide12_df.table.rows:
                    row.height = Cm(0.8)
    filename = datetime.now().strftime("%Y%m%d-%H%M%S")
    presentation.save(path_to_pptx + '/' + story_title + '_' + filename + '.pptx')

    # # Default slide = 3 - empty slide with header
    #
    #
    # def analyze_ppt(input, output):
    #     """ Take the input file and analyze the structure.
    #     The output file contains marked up information to make it easier
    #     for generating future powerpoint templates.
    #     """
    #     prs = Presentation(input)
    #     # Each powerpoint file has multiple layouts
    #     # Loop through them all and  see where the various elements are
    #     for index, _ in enumerate(prs.slide_layouts):
    #         slide = prs.slides.add_slide(prs.slide_layouts[index])
    #         # Not every slide has to have a title
    #         try:
    #             title = slide.shapes.title
    #             title.text = 'Title for Layout {}'.format(index)
    #         except AttributeError:
    #             print("No Title for Layout {}".format(index))
    #         # Go through all the placeholders and identify them by index and type
    #         for shape in slide.placeholders:
    #             if shape.is_placeholder:
    #                 phf = shape.placeholder_format
    #                 # Do not overwrite the title which is just a special placeholder
    #                 try:
    #                     if 'Title' not in shape.text:
    #                         shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
    #                 except AttributeError:
    #                     print("{} has no text attribute".format(phf.type))
    #                 print('{} {}'.format(phf.idx, shape.name))
    #     prs.save(output)
    #
    # prs_path = '/Users/annademidova/PycharmProjects/HAR_Parser/'
    # analyze_ppt(settings.TEMPLATE_PATH, prs_path + 'Analysis.pptx')